# build_course_package.py
# Generates: PDF + DOCX + PPTX with branding, flow diagrams (Loops/DS/StudentDB),
# practice exercises+answers after each topic, and MCQs at end of each chapter.

import os, datetime, textwrap, copy, json
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, FancyArrowPatch

# PDF (reportlab)
from reportlab.lib.pagesizes import A4
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, PageBreak, Preformatted, Image, Table, TableStyle
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib import colors
from reportlab.pdfgen import canvas

# DOCX
from docx import Document
from docx.shared import Inches

# PPTX
from pptx import Presentation
from pptx.util import Inches as PInches


# -------------------------
# CONFIG / BRANDING
# -------------------------
BRAND_LINE_1 = "Programmer’s Picnic"
BRAND_LINE_2 = "Python Starter Course — 7 PM Batch"
AUTHOR = "Champak Roy"

OUT_DIR = os.path.abspath(".")
DATE_STR = datetime.date.today().strftime("%d %b %Y")

PDF_PATH  = os.path.join(OUT_DIR, "Python_Starter_Course_FULL_Package.pdf")
DOCX_PATH = os.path.join(OUT_DIR, "Python_Starter_Course_FULL_Package.docx")
PPTX_PATH = os.path.join(OUT_DIR, "Python_Starter_Course_FULL_Package.pptx")


# -------------------------
# FLOW DIAGRAMS (requested)
# Loops: 6–10
# DS: 11–15
# StudentDB: 18–19
# -------------------------
DIAGRAM_TOPIC_IDS = [6,7,8,9,10,11,12,13,14,15,18,19]

def save_flow_diagram(title, nodes, edges, notes, path_png, figsize=(8,2.8)):
    fig, ax = plt.subplots(figsize=figsize)
    ax.axis("off")
    ax.text(0.01, 0.93, title, fontsize=14, weight="bold", transform=ax.transAxes)
    for (x,y,w,h,label) in nodes:
        ax.add_patch(Rectangle((x,y), w,h, fill=False, linewidth=1.5, transform=ax.transAxes))
        ax.text(x+w/2, y+h/2, label, ha="center", va="center", fontsize=10, transform=ax.transAxes)
    for (x1,y1,x2,y2) in edges:
        ax.add_patch(FancyArrowPatch((x1,y1),(x2,y2), arrowstyle="->",
                                     mutation_scale=12, linewidth=1.5, transform=ax.transAxes))
    if notes:
        ax.text(0.01, 0.06, notes, fontsize=9.5, transform=ax.transAxes)
    plt.tight_layout()
    plt.savefig(path_png, dpi=180)
    plt.close(fig)

def diagram_for_topic(tid:int) -> str:
    # We reuse StudentDB diagram for both 18 and 19
    if tid == 19:
        tid = 18
    path = os.path.join(OUT_DIR, f"diag_topic_{tid:02d}.png")

    if tid in (6,8):  # while
        save_flow_diagram(
            "while Loop Flow",
            [(0.05,0.62,0.30,0.22,"Start/Init"),
             (0.40,0.62,0.30,0.22,"Condition?"),
             (0.40,0.34,0.30,0.22,"Body+Update"),
             (0.75,0.48,0.20,0.22,"End")],
            [(0.35,0.73,0.40,0.73),(0.55,0.62,0.55,0.56),
             (0.55,0.34,0.55,0.62),(0.70,0.73,0.75,0.59)],
            "Update variables to avoid infinite loop.", path, figsize=(8,3.0)
        )
    elif tid == 7:  # for range
        save_flow_diagram(
            "for Loop with range()",
            [(0.05,0.60,0.30,0.24,"range(...)"),
             (0.40,0.60,0.30,0.24,"i takes values"),
             (0.75,0.60,0.20,0.24,"Body")],
            [(0.35,0.72,0.40,0.72),(0.70,0.72,0.75,0.72)],
            "Best when count is known.", path, figsize=(8,2.6)
        )
    elif tid in (9,10):  # nested
        save_flow_diagram(
            "Nested Loops",
            [(0.05,0.58,0.28,0.26,"Outer loop\n(rows)"),
             (0.38,0.58,0.28,0.26,"Inner loop\n(cols)"),
             (0.71,0.58,0.24,0.26,"print/check")],
            [(0.33,0.71,0.38,0.71),(0.66,0.71,0.71,0.71)],
            "Outer controls lines; inner controls items per line.", path, figsize=(8,2.6)
        )
    elif tid == 11:  # list
        save_flow_diagram(
            "List (Ordered, Mutable)",
            [(0.05,0.58,0.27,0.26,"Create\n[]"),
             (0.36,0.58,0.27,0.26,"append/remove"),
             (0.67,0.58,0.28,0.26,"index/slice")],
            [(0.32,0.71,0.36,0.71),(0.63,0.71,0.67,0.71)],
            "Sequence where order matters.", path, figsize=(8,2.6)
        )
    elif tid == 12:  # dict
        save_flow_diagram(
            "Dictionary (Key → Value)",
            [(0.05,0.58,0.30,0.26,"key"),
             (0.38,0.58,0.30,0.26,"maps to"),
             (0.71,0.58,0.24,0.26,"value")],
            [(0.35,0.71,0.38,0.71),(0.68,0.71,0.71,0.71)],
            "db[roll] → record.", path, figsize=(8,2.6)
        )
    elif tid == 13:  # tuple
        save_flow_diagram(
            "Tuple (Ordered, Immutable)",
            [(0.05,0.58,0.30,0.26,"Create\n( )"),
             (0.38,0.58,0.30,0.26,"unpack"),
             (0.71,0.58,0.24,0.26,"use safely")],
            [(0.35,0.71,0.38,0.71),(0.68,0.71,0.71,0.71)],
            "Use when values should not change.", path, figsize=(8,2.6)
        )
    elif tid == 14:  # set
        save_flow_diagram(
            "Set (Unique)",
            [(0.05,0.58,0.30,0.26,"Create\n{ }"),
             (0.38,0.58,0.30,0.26,"add/remove"),
             (0.71,0.58,0.24,0.26,"ops | & -")],
            [(0.35,0.71,0.38,0.71),(0.68,0.71,0.71,0.71)],
            "Duplicates removed automatically.", path, figsize=(8,2.6)
        )
    elif tid == 15:  # frozenset
        save_flow_diagram(
            "frozenset (Immutable Set)",
            [(0.05,0.58,0.30,0.26,"frozenset(...)"),
             (0.38,0.58,0.30,0.26,"no change"),
             (0.71,0.58,0.24,0.26,"dict key")],
            [(0.35,0.71,0.38,0.71),(0.68,0.71,0.71,0.71)],
            "Hashable and safe.", path, figsize=(8,2.6)
        )
    else:  # StudentDB architecture
        save_flow_diagram(
            "StudentDB Architecture",
            [(0.04,0.62,0.22,0.22,"Menu"),
             (0.30,0.62,0.22,0.22,"CRUD"),
             (0.56,0.62,0.22,0.22,"Undo/Redo"),
             (0.80,0.62,0.18,0.22,"Save/Load")],
            [(0.26,0.73,0.30,0.73),(0.52,0.73,0.56,0.73),(0.78,0.73,0.80,0.73)],
            "Undo stores snapshots; redo cleared on new change.", path, figsize=(8,3.0)
        )

    return path


# -------------------------
# YOUR 19 TOPICS (brief but complete)
# Each topic: intro/explain/importance + 3 code samples
# -------------------------
TOPICS = [
  (1,"Introduction to Python, Installation, Hello World",
   "Python is a beginner-friendly language used in automation, web, data, and scripting.",
   "Install Python, install VS Code, add Python extension, create a .py file and run it.",
   "A working setup is the base of every project.",
   ['print("Hello World")',
    'print("Welcome to Python!")\nprint("Programmer\'s Picnic - 7 PM Batch")',
    'course="Python Starter"\nprint("Course:", course)']),
  (2,"Input, print and f-strings",
   "Input makes programs interactive; output shows results.",
   "input() returns str; convert using int()/float(). f-strings format text as f\"...{x}...\".",
   "Used in menus/projects: data entry, reports, receipts.",
   ['name=input("Enter name: ")\nprint("Hello", name)',
    'age=int(input("Enter age: "))\nprint(f"Next year: {age+1}")',
    'a=int(input("A: ")); b=int(input("B: "))\nprint(f"{a}+{b}={a+b}")']),
  (3,"Arithmetic Operators + BODMAS",
   "Arithmetic operators power calculators and scoring systems.",
   "+ - * / // % ** ; precedence: () > ** > * / // % > + -",
   "Correct precedence avoids wrong answers; brackets give control.",
   ["a,b=10,3\nprint(a+b,a-b,a*b)\nprint(a/b,a//b,a%b)\nprint(a**b)",
    "print(10 + 2 * 3)\nprint((10 + 2) * 3)",
    "total=251\nstudents=4\nprint('Average:', total/students)\nprint('Remainder:', total%students)"]),
  (4,"Conditional + Logical Operators",
   "Conditions allow decisions based on data.",
   "Relational: > < >= <= == != ; Logical: and, or, not.",
   "Used in eligibility, validation, grading, filtering.",
   ["x=12\nprint(x>10)\nprint(x==12)\nprint(x!=5)",
    "age=19\nhas_id=True\nprint(age>=18 and has_id)\nprint(age<18 or not has_id)",
    "marks=72\nattendance=80\nprint(marks>=33 and attendance>=75)"]),
  (5,"If-Else + Ternary",
   "if/elif/else chooses one path among many.",
   "Use if/elif/else blocks; ternary: a if condition else b.",
   "Core logic of every app: choices, menus, rules.",
   ["n=int(input('Enter n: '))\nprint('Even' if n%2==0 else 'Odd')",
    "m=int(input('Marks: '))\nprint('Pass' if m>=33 else 'Fail')",
    "x=int(input('x: '))\nif x>0:\n    print('Positive')\nelse:\n    print('Non-positive')"]),
  (6,"Loops: while (intro)",
   "Loops repeat tasks automatically.",
   "while runs while condition is True; update variables to avoid infinite loops.",
   "Used for counters, repeated input, menus.",
   ["i=1\nwhile i<=5:\n    print(i)\n    i+=1",
    "s=0\ni=1\nwhile i<=5:\n    s+=i\n    i+=1\nprint('Sum:', s)",
    "while True:\n    if input('Type q to quit: ')== 'q':\n        break"]),
  (7,"for loops using range",
   "for loops iterate over counts and sequences.",
   "range(stop), range(start, stop), range(start, stop, step).",
   "Used in tables, patterns, list processing.",
   ["for i in range(5):\n    print(i)",
    "for i in range(1,6):\n    print(i)",
    "for i in range(10,0,-2):\n    print(i)"]),
  (8,"while loops (practice)",
   "While loops are perfect for user-driven repetition.",
   "Use break/continue; validation and retries.",
   "Password retries, validation loops, menus.",
   ["attempts=0\nwhile attempts<3:\n    attempts+=1\n    print('Try', attempts)",
    "total=0\nwhile True:\n    n=int(input('n (0 stop): '))\n    if n==0: break\n    total+=n\nprint(total)",
    "i=0\nwhile i<10:\n    i+=1\n    if i%2==0: continue\n    print(i)"]),
  (9,"Nested loops: prime numbers",
   "Nested loops help in factor checking.",
   "Prime numbers have exactly two factors; optimize to sqrt(n).",
   "Builds strong logic + efficiency thinking.",
   ["n=29\nprime=True\nfor i in range(2,int(n**0.5)+1):\n    if n%i==0:\n        prime=False\nprint(prime)",
    "n=int(input('n: '))\nprime=True\nif n<2: prime=False\nfor i in range(2,int(n**0.5)+1):\n    if n%i==0:\n        prime=False; break\nprint('Prime' if prime else 'Not prime')",
    "for n in range(2,51):\n    prime=True\n    for i in range(2,int(n**0.5)+1):\n        if n%i==0:\n            prime=False; break\n    if prime:\n        print(n,end=' ')\nprint()"]),
  (10,"Nested loops: patterns",
   "Patterns teach row/column thinking.",
   "Outer loop controls rows; inner prints columns; use end='' for same line.",
   "Fastest way to master nested loops.",
   ["for r in range(1,6):\n    print('*'*r)",
    "for r in range(5,0,-1):\n    print('*'*r)",
    "for r in range(1,6):\n    print(' '.join(str(x) for x in range(1,r+1)))"]),
  (11,"List",
   "Lists store ordered, changeable sequences.",
   "Indexing, slicing, append/remove, sort, loop.",
   "Used for collections, results, history.",
   ["a=[10,20,30]\nprint(a[0],a[-1])",
    "a=[1,2,3]\na.append(4)\na.remove(2)\nprint(a)",
    "nums=[5,2,9,1]\nnums.sort()\nprint(nums)"]),
  (12,"Dictionary",
   "Dictionaries store key-value records for fast lookup.",
   "Use dicts for StudentDB: roll → record; iterate with items().",
   "Represents real-world entities: students, users, products.",
   ["student={'roll':1,'name':'Amit'}\nprint(student['name'])",
    "student={'roll':1,'name':'Amit'}\nstudent['marks']=88\nprint(student)",
    "db={101:{'name':'Ravi'},102:{'name':'Neha'}}\nfor roll,info in db.items():\n    print(roll, info['name'])"]),
  (13,"Tuple",
   "Tuples store ordered but immutable data.",
   "Useful for fixed values and unpacking.",
   "Safer than lists when data should not change.",
   ["t=(10,20)\nprint(t)",
    "x,y=(5,7)\nprint(x,y)",
    "coords=[(1,2),(3,4)]\nfor x,y in coords:\n    print(x+y)"]),
  (14,"Set",
   "Sets store unique elements (no duplicates).",
   "Operations: union(|), intersection(&), difference(-).",
   "Great for de-duplication and comparisons.",
   ["s={1,2,2,3}\nprint(s)",
    "a={1,2,3}; b={3,4}\nprint(a|b)",
    "a={1,2,3}; b={2,3,5}\nprint(a&b)\nprint(a-b)"]),
  (15,"Frozen DS (frozenset)",
   "frozenset is an immutable set.",
   "You can’t change it; it’s hashable (can be dict key).",
   "Useful when you need safe, fixed groups.",
   ["fs=frozenset([1,2,2,3])\nprint(fs)",
    "group=frozenset(['A','B'])\nprint('A' in group)",
    "d={frozenset([1,2]):'pair'}\nprint(d[frozenset([1,2])])"]),
  (16,"Menu using infinite while loops",
   "Menu-driven programs run until Exit is chosen.",
   "while True + if/elif; break exits; each option calls a function.",
   "Pattern powers StudentDB and most console apps.",
   ["while True:\n    print('1.Hi 2.Exit')\n    ch=input('Choice: ')\n    if ch=='1': print('Hi!')\n    elif ch=='2': break",
    "items=[]\nwhile True:\n    print('1.Add 2.View 3.Exit')\n    ch=input('Choice: ')\n    if ch=='1': items.append(input('Item: '))\n    elif ch=='2': print(items)\n    elif ch=='3': break",
    "while True:\n    ch=input('Enter y/n: ')\n    if ch in ('y','n'):\n        print('OK'); break\n    print('Invalid')"]),
  (17,"Pickle",
   "Pickle saves Python objects to a binary file.",
   "pickle.dump(obj, file) to save; pickle.load(file) to load; use try/except.",
   "Fast Python-only persistence (not human-readable).",
   ["import pickle\npickle.dump({'a':1}, open('x.pkl','wb'))",
    "import pickle\nprint(pickle.load(open('x.pkl','rb')))",
    "import os\nprint(os.path.exists('x.pkl'))"]),
  (18,"StudentDB using Dictionary + Pickle (Undo/Redo)",
   "A mini database using dict + pickle persistence.",
   "CRUD + save/load + undo/redo snapshots for safety.",
   "Your first complete console project: data + logic + persistence.",
   ["# Full code in Appendix A",
    "# Run: python studentdb_pickle_undo_redo.py",
    "# Try: add → undo → redo → save → exit → run again → load"]),
  (19,"StudentDB using Dictionary + JSON (Undo/Redo)",
   "Same StudentDB using JSON (human-readable).",
   "Store keys as strings; json.dump with indent; same undo/redo logic.",
   "Great for portability and debugging saved data.",
   ["# Full code in Appendix B",
    "# Run: python studentdb_json_undo_redo.py",
    "# Open studentdb.json to view records"]),
]

# -------------------------
# EXERCISES + ANSWERS (3 each)
# -------------------------
def exercises_for(topic_title:str):
    ex = [
        f"Exercise 1: Write a small program related to: {topic_title}.",
        "Exercise 2: Modify one code sample to take user input and show output using f-strings.",
        "Exercise 3: Add validation (reject wrong input) and test with 3 cases."
    ]
    ans = [
        "Answer 1: Implement core concept + test with 2–3 runs.",
        "Answer 2: Use input(), int()/float() conversions, and f-strings.",
        "Answer 3: Use while loop + if checks; re-ask until valid."
    ]
    return ex, ans


# -------------------------
# CHAPTERS + MCQs
# -------------------------
CHAPTERS = [
    ("Chapter 1: Basics & I/O", [1,2]),
    ("Chapter 2: Operators & Decisions", [3,4,5]),
    ("Chapter 3: Loops & Nested Loops", [6,7,8,9,10]),
    ("Chapter 4: Data Structures", [11,12,13,14,15]),
    ("Chapter 5: Persistence & Projects", [16,17,18,19]),
]

MCQS = {
 "Chapter 1: Basics & I/O":[
  ("input() returns what type?", ["A) int","B) str","C) float","D) list"], "B"),
  ("Python file extension?", ["A) .py","B) .js","C) .html","D) .txt"], "A"),
  ("f-string starts with:", ["A) g","B) f","C) p","D) x"], "B"),
  ("VS Code is:", ["A) compiler","B) editor","C) OS","D) database"], "B"),
  ("print() does:", ["A) input","B) output","C) save","D) delete"], "B"),
 ],
 "Chapter 2: Operators & Decisions":[
  ("Highest precedence:", ["A) +","B) *","C) **","D) -"], "C"),
  ("10//3 equals:", ["A) 3","B) 3.33","C) 0","D) 1"], "A"),
  ("== means:", ["A) assign","B) equal","C) not","D) type"], "B"),
  ("and true when:", ["A) any true","B) both true","C) none","D) random"], "B"),
  ("Ternary form:", ["A) a if c else b","B) if a else b","C) c if a else b","D) a else c"], "A"),
 ],
 "Chapter 3: Loops & Nested Loops":[
  ("while runs while:", ["A) once","B) condition true","C) always false","D) compiled"], "B"),
  ("range(5) gives:", ["A) 1..5","B) 0..4","C) 0..5","D) 5..0"], "B"),
  ("break does:", ["A) skip","B) exit loop","C) restart","D) pause"], "B"),
  ("Nested loops used for:", ["A) patterns","B) imports","C) errors","D) comments"], "A"),
  ("Prime optimization uses:", ["A) n","B) n/2","C) sqrt(n)","D) 1"], "C"),
 ],
 "Chapter 4: Data Structures":[
  ("List is:", ["A) mutable","B) immutable","C) key-value","D) unique only"], "A"),
  ("Dict uses:", ["A) indexes","B) keys","C) slices","D) loops"], "B"),
  ("Tuple is:", ["A) mutable","B) immutable","C) unordered","D) unique"], "B"),
  ("Set stores:", ["A) duplicates","B) unique elements","C) key-values","D) order"], "B"),
  ("frozenset is:", ["A) mutable","B) immutable","C) list","D) dict"], "B"),
 ],
 "Chapter 5: Persistence & Projects":[
  ("Menu loops often use:", ["A) while True","B) if only","C) recursion","D) compile"], "A"),
  ("Pickle stores:", ["A) text","B) binary","C) images","D) audio"], "B"),
  ("JSON is:", ["A) binary","B) human-readable text","C) python-only","D) encrypted"], "B"),
  ("Undo/Redo uses:", ["A) stacks","B) random","C) files only","D) sleep"], "A"),
  ("Redo cleared when:", ["A) undo","B) new change","C) save","D) exit"], "B"),
 ],
}

# -------------------------
# FULL PROJECTS (Undo/Redo included)
# -------------------------
STUDENTDB_PICKLE_CODE = r'''# studentdb_pickle_undo_redo.py
import pickle
import os
import copy

FILE_NAME = "studentdb.pkl"

def load_db():
    if os.path.exists(FILE_NAME):
        try:
            with open(FILE_NAME, "rb") as f:
                db = pickle.load(f)
            if isinstance(db, dict):
                cleaned = {}
                for k, v in db.items():
                    try:
                        ik = int(k)
                    except Exception:
                        continue
                    cleaned[ik] = v
                return cleaned
        except Exception:
            pass
    return {}

def save_db(db):
    with open(FILE_NAME, "wb") as f:
        pickle.dump(db, f)

def input_int(prompt):
    while True:
        s = input(prompt).strip()
        if s.isdigit():
            return int(s)
        print("Please enter a valid number.")

def snapshot(db):
    return copy.deepcopy(db)

def push_undo(undo_stack, redo_stack, db):
    undo_stack.append(snapshot(db))
    redo_stack.clear()

def add_student(db, undo_stack, redo_stack):
    push_undo(undo_stack, redo_stack, db)
    roll = input_int("Enter roll: ")
    if roll in db:
        print("Roll already exists. Use update option.")
        return
    name = input("Enter name: ").strip()
    marks = input_int("Enter marks (0-100): ")
    db[roll] = {"roll": roll, "name": name, "marks": marks}
    print("Student added.")

def view_all(db):
    if not db:
        print("No records found.")
        return
    print("\n--- All Students ---")
    for roll in sorted(db.keys()):
        s = db[roll]
        print(f"Roll: {s['roll']} | Name: {s['name']} | Marks: {s['marks']}")

def search_student(db):
    roll = input_int("Enter roll to search: ")
    s = db.get(roll)
    if not s:
        print("Student not found.")
        return
    print(f"Found -> Roll: {s['roll']} | Name: {s['name']} | Marks: {s['marks']}")

def update_student(db, undo_stack, redo_stack):
    roll = input_int("Enter roll to update: ")
    if roll not in db:
        print("Student not found.")
        return
    push_undo(undo_stack, redo_stack, db)
    s = db[roll]
    print(f"Current -> Name: {s['name']} | Marks: {s['marks']}")
    new_name = input("Enter new name (blank = keep): ").strip()
    new_marks = input("Enter new marks (blank = keep): ").strip()
    if new_name:
        s["name"] = new_name
    if new_marks:
        if new_marks.isdigit():
            s["marks"] = int(new_marks)
        else:
            print("Invalid marks. Keeping old marks.")
    db[roll] = s
    print("Student updated.")

def delete_student(db, undo_stack, redo_stack):
    roll = input_int("Enter roll to delete: ")
    if roll not in db:
        print("Student not found.")
        return
    push_undo(undo_stack, redo_stack, db)
    del db[roll]
    print("Student deleted.")

def do_undo(db, undo_stack, redo_stack):
    if not undo_stack:
        print("Nothing to undo.")
        return db
    redo_stack.append(snapshot(db))
    prev = undo_stack.pop()
    print("Undo done.")
    return prev

def do_redo(db, undo_stack, redo_stack):
    if not redo_stack:
        print("Nothing to redo.")
        return db
    undo_stack.append(snapshot(db))
    nxt = redo_stack.pop()
    print("Redo done.")
    return nxt

def main():
    db = load_db()
    undo_stack = []
    redo_stack = []
    print("StudentDB (Pickle) loaded. Records:", len(db))

    while True:
        print("""
=========================
 StudentDB (Pickle)
=========================
1. Add Student
2. View All Students
3. Search Student
4. Update Student
5. Delete Student
6. Undo
7. Redo
8. Save
9. Save & Exit
10. Exit (No Save)
""")
        choice = input("Choose option: ").strip()

        if choice == "1":
            add_student(db, undo_stack, redo_stack)
        elif choice == "2":
            view_all(db)
        elif choice == "3":
            search_student(db)
        elif choice == "4":
            update_student(db, undo_stack, redo_stack)
        elif choice == "5":
            delete_student(db, undo_stack, redo_stack)
        elif choice == "6":
            db = do_undo(db, undo_stack, redo_stack)
        elif choice == "7":
            db = do_redo(db, undo_stack, redo_stack)
        elif choice == "8":
            save_db(db)
            print("Saved.")
        elif choice == "9":
            save_db(db)
            print("Saved. Bye!")
            break
        elif choice == "10":
            print("Bye! (Not saved)")
            break
        else:
            print("Invalid choice.")

if __name__ == "__main__":
    main()
'''

STUDENTDB_JSON_CODE = r'''# studentdb_json_undo_redo.py
import json
import os
import copy

FILE_NAME = "studentdb.json"

def load_db():
    if os.path.exists(FILE_NAME):
        try:
            with open(FILE_NAME, "r", encoding="utf-8") as f:
                db = json.load(f)
            if isinstance(db, dict):
                cleaned = {}
                for k, v in db.items():
                    if isinstance(k, str) and isinstance(v, dict):
                        cleaned[k] = v
                return cleaned
        except Exception:
            pass
    return {}

def save_db(db):
    with open(FILE_NAME, "w", encoding="utf-8") as f:
        json.dump(db, f, indent=2, ensure_ascii=False)

def input_int(prompt):
    while True:
        s = input(prompt).strip()
        if s.isdigit():
            return int(s)
        print("Please enter a valid number.")

def snapshot(db):
    return copy.deepcopy(db)

def push_undo(undo_stack, redo_stack, db):
    undo_stack.append(snapshot(db))
    redo_stack.clear()

def add_student(db, undo_stack, redo_stack):
    push_undo(undo_stack, redo_stack, db)
    roll = input_int("Enter roll: ")
    key = str(roll)
    if key in db:
        print("Roll already exists. Use update option.")
        return
    name = input("Enter name: ").strip()
    marks = input_int("Enter marks (0-100): ")
    db[key] = {"roll": roll, "name": name, "marks": marks}
    print("Student added.")

def view_all(db):
    if not db:
        print("No records found.")
        return
    print("\n--- All Students ---")
    for key in sorted(db.keys(), key=lambda x: int(x)):
        s = db[key]
        print(f"Roll: {s['roll']} | Name: {s['name']} | Marks: {s['marks']}")

def search_student(db):
    roll = input_int("Enter roll to search: ")
    s = db.get(str(roll))
    if not s:
        print("Student not found.")
        return
    print(f"Found -> Roll: {s['roll']} | Name: {s['name']} | Marks: {s['marks']}")

def update_student(db, undo_stack, redo_stack):
    roll = input_int("Enter roll to update: ")
    key = str(roll)
    if key not in db:
        print("Student not found.")
        return
    push_undo(undo_stack, redo_stack, db)
    s = db[key]
    print(f"Current -> Name: {s['name']} | Marks: {s['marks']}")
    new_name = input("Enter new name (blank = keep): ").strip()
    new_marks = input("Enter new marks (blank = keep): ").strip()
    if new_name:
        s["name"] = new_name
    if new_marks:
        if new_marks.isdigit():
            s["marks"] = int(new_marks)
        else:
            print("Invalid marks. Keeping old marks.")
    db[key] = s
    print("Student updated.")

def delete_student(db, undo_stack, redo_stack):
    roll = input_int("Enter roll to delete: ")
    key = str(roll)
    if key not in db:
        print("Student not found.")
        return
    push_undo(undo_stack, redo_stack, db)
    del db[key]
    print("Student deleted.")

def do_undo(db, undo_stack, redo_stack):
    if not undo_stack:
        print("Nothing to undo.")
        return db
    redo_stack.append(snapshot(db))
    prev = undo_stack.pop()
    print("Undo done.")
    return prev

def do_redo(db, undo_stack, redo_stack):
    if not redo_stack:
        print("Nothing to redo.")
        return db
    undo_stack.append(snapshot(db))
    nxt = redo_stack.pop()
    print("Redo done.")
    return nxt

def main():
    db = load_db()
    undo_stack = []
    redo_stack = []
    print("StudentDB (JSON) loaded. Records:", len(db))

    while True:
        print("""
=========================
 StudentDB (JSON)
=========================
1. Add Student
2. View All Students
3. Search Student
4. Update Student
5. Delete Student
6. Undo
7. Redo
8. Save
9. Save & Exit
10. Exit (No Save)
""")
        choice = input("Choose option: ").strip()

        if choice == "1":
            add_student(db, undo_stack, redo_stack)
        elif choice == "2":
            view_all(db)
        elif choice == "3":
            search_student(db)
        elif choice == "4":
            update_student(db, undo_stack, redo_stack)
        elif choice == "5":
            delete_student(db, undo_stack, redo_stack)
        elif choice == "6":
            db = do_undo(db, undo_stack, redo_stack)
        elif choice == "7":
            db = do_redo(db, undo_stack, redo_stack)
        elif choice == "8":
            save_db(db)
            print("Saved.")
        elif choice == "9":
            save_db(db)
            print("Saved. Bye!")
            break
        elif choice == "10":
            print("Bye! (Not saved)")
            break
        else:
            print("Invalid choice.")

if __name__ == "__main__":
    main()
'''

# -------------------------
# BUILDERS
# -------------------------
def pdf_footer(c: canvas.Canvas, doc):
    c.saveState()
    c.setFont("Helvetica", 9)
    c.setFillColor(colors.grey)
    c.drawString(1.5*cm, 1.1*cm, f"{BRAND_LINE_1} — {BRAND_LINE_2}")
    c.drawRightString(A4[0]-1.5*cm, 1.1*cm, f"Page {c.getPageNumber()}")
    c.restoreState()

def code_block(styles, txt):
    wrap=[]
    for line in txt.splitlines():
        wrap += textwrap.wrap(line, 95) if len(line) > 95 else [line]
    return Preformatted("\n".join(wrap), styles["Code"])

def build_pdf(diag_paths):
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="T", parent=styles["Title"], fontSize=26, leading=30))
    styles.add(ParagraphStyle(name="ST", parent=styles["Heading2"], fontSize=14, leading=18,
                              textColor=colors.HexColor("#d97706")))
    styles.add(ParagraphStyle(name="H1", parent=styles["Heading1"], fontSize=18, leading=22,
                              textColor=colors.HexColor("#d97706")))
    styles.add(ParagraphStyle(name="H2", parent=styles["Heading2"], fontSize=14, leading=18))
    styles.add(ParagraphStyle(name="B", parent=styles["BodyText"], fontSize=10.5, leading=14))
    styles.add(ParagraphStyle(name="Code", parent=styles["BodyText"], fontName="Courier",
                              fontSize=8.4, leading=10))

    story=[]
    story.append(Paragraph("Python Starter Course", styles["T"]))
    story.append(Paragraph(BRAND_LINE_1, styles["ST"]))
    story.append(Paragraph(BRAND_LINE_2, styles["ST"]))
    story.append(Spacer(1,8))
    story.append(Paragraph(f"<b>Prepared by:</b> {AUTHOR} &nbsp;&nbsp;|&nbsp;&nbsp; <b>Date:</b> {DATE_STR}", styles["B"]))
    story.append(Spacer(1,10))
    story.append(Paragraph(
        "This edition includes: flow diagrams (Loops 6–10, Data Structures 11–15, StudentDB 18–19), "
        "practice exercises with answers after every topic, and MCQs at the end of each chapter.",
        styles["B"]
    ))
    story.append(PageBreak())

    # Contents table
    story.append(Paragraph("Contents", styles["H1"]))
    toc=[["Chapter","Topics"]]+[[ch, ", ".join(str(i) for i in ids)] for ch,ids in CHAPTERS]
    tbl=Table(toc, colWidths=[6*cm, 10.5*cm])
    tbl.setStyle(TableStyle([
        ("GRID",(0,0),(-1,-1),0.5,colors.grey),
        ("BACKGROUND",(0,0),(-1,0),colors.lightgrey),
        ("FONTNAME",(0,0),(-1,0),"Helvetica-Bold")
    ]))
    story.append(tbl)
    story.append(PageBreak())

    topic_map = {t[0]:t for t in TOPICS}

    for ch, ids in CHAPTERS:
        story.append(Paragraph(ch, styles["H1"]))
        story.append(Paragraph(f"Topics: {', '.join(str(i) for i in ids)}", styles["B"]))
        story.append(Spacer(1,6))

        for tid in ids:
            _, title, intro, expl, imp, samples = topic_map[tid]
            story.append(Paragraph(f"{tid}. {title}", styles["H2"]))
            story.append(Paragraph(f"<b>Introduction:</b> {intro}", styles["B"]))
            story.append(Paragraph(f"<b>Explanation:</b> {expl}", styles["B"]))
            story.append(Paragraph(f"<b>Importance:</b> {imp}", styles["B"]))
            story.append(Spacer(1,6))

            if tid in DIAGRAM_TOPIC_IDS:
                story.append(Paragraph("<b>Flow Diagram</b>", styles["B"]))
                img_id = 18 if tid==19 else tid
                story.append(Image(diag_paths[img_id], width=17.0*cm, height=5.2*cm))
                story.append(Spacer(1,6))

            story.append(Paragraph("<b>Code Samples (3)</b>", styles["B"]))
            for i, s in enumerate(samples, 1):
                story.append(Paragraph(f"Sample {i}", styles["B"]))
                story.append(code_block(styles, s))
                story.append(Spacer(1,4))

            ex, ans = exercises_for(title)
            story.append(Paragraph("<b>Practice Exercises</b>", styles["B"]))
            for line in ex:
                story.append(Paragraph("• " + line, styles["B"]))
            story.append(Spacer(1,4))
            story.append(Paragraph("<b>Answers</b>", styles["B"]))
            for line in ans:
                story.append(Paragraph("• " + line, styles["B"]))
            story.append(Spacer(1,10))

        # MCQs
        story.append(Paragraph("MCQs (End of Chapter)", styles["H2"]))
        for i,(q,opts,key) in enumerate(MCQS[ch],1):
            story.append(Paragraph(f"<b>Q{i}.</b> {q}", styles["B"]))
            for o in opts:
                story.append(Paragraph("&nbsp;&nbsp;" + o, styles["B"]))
        story.append(Spacer(1,6))
        story.append(Paragraph("<b>Answer Key:</b> " + ", ".join([f"Q{i+1}-{MCQS[ch][i][2]}" for i in range(len(MCQS[ch]))]), styles["B"]))
        story.append(PageBreak())

    # Appendices: full projects
    story.append(Paragraph("Appendix A — StudentDB Pickle Project (Full Code)", styles["H1"]))
    story.append(code_block(styles, STUDENTDB_PICKLE_CODE))
    story.append(PageBreak())
    story.append(Paragraph("Appendix B — StudentDB JSON Project (Full Code)", styles["H1"]))
    story.append(code_block(styles, STUDENTDB_JSON_CODE))

    doc = SimpleDocTemplate(PDF_PATH, pagesize=A4, leftMargin=1.6*cm, rightMargin=1.6*cm,
                            topMargin=1.6*cm, bottomMargin=1.6*cm)
    doc.build(story, onFirstPage=pdf_footer, onLaterPages=pdf_footer)

def build_docx(diag_paths):
    d = Document()
    d.add_heading("Python Starter Course", level=0)
    d.add_paragraph(BRAND_LINE_1)
    d.add_paragraph(BRAND_LINE_2)
    d.add_paragraph(f"Prepared by {AUTHOR} | {DATE_STR}")
    d.add_page_break()

    topic_map = {t[0]:t for t in TOPICS}

    for ch, ids in CHAPTERS:
        d.add_heading(ch, level=1)
        d.add_paragraph("Topics: " + ", ".join(str(i) for i in ids))

        for tid in ids:
            _, title, intro, expl, imp, samples = topic_map[tid]
            d.add_heading(f"{tid}. {title}", level=2)
            d.add_paragraph("Introduction: " + intro)
            d.add_paragraph("Explanation: " + expl)
            d.add_paragraph("Importance: " + imp)

            if tid in DIAGRAM_TOPIC_IDS:
                img_id = 18 if tid==19 else tid
                d.add_paragraph("Flow Diagram:")
                d.add_picture(diag_paths[img_id], width=Inches(6.5))

            d.add_paragraph("Code Samples:")
            for s in samples:
                d.add_paragraph(s)

            ex, ans = exercises_for(title)
            d.add_paragraph("Practice Exercises:")
            for line in ex:
                d.add_paragraph(line)
            d.add_paragraph("Answers:")
            for line in ans:
                d.add_paragraph(line)

        d.add_heading("MCQs (End of Chapter)", level=2)
        for i,(q,opts,key) in enumerate(MCQS[ch],1):
            d.add_paragraph(f"Q{i}. {q}")
            for o in opts:
                d.add_paragraph("   " + o)
        d.add_paragraph("Answer Key: " + ", ".join([f"Q{i+1}-{MCQS[ch][i][2]}" for i in range(len(MCQS[ch]))]))
        d.add_page_break()

    d.add_heading("Appendix — Full Projects", level=1)
    d.add_paragraph("Appendix A: studentdb_pickle_undo_redo.py (full code included in PDF).")
    d.add_paragraph("Appendix B: studentdb_json_undo_redo.py (full code included in PDF).")
    d.save(DOCX_PATH)

def build_ppt(diag_paths):
    prs = Presentation()

    # Title slide
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = "Python Starter Course"
    s0.placeholders[1].text = f"{BRAND_LINE_1}\n{BRAND_LINE_2}\n{AUTHOR}\n{DATE_STR}"

    # Chapter slides
    for ch, ids in CHAPTERS:
        overview = prs.slides.add_slide(prs.slide_layouts[1])
        overview.shapes.title.text = ch
        tf = overview.placeholders[1].text_frame
        tf.clear()
        tf.text = "Topics: " + ", ".join(str(i) for i in ids)

        for tid in ids:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
            slide.shapes.title.text = f"{tid}. {TOPICS[tid-1][1]}"

            if tid in DIAGRAM_TOPIC_IDS:
                img_id = 18 if tid==19 else tid
                slide.shapes.add_picture(diag_paths[img_id], PInches(0.8), PInches(1.6), width=PInches(8.3))
            else:
                box = slide.shapes.add_textbox(PInches(0.8), PInches(1.6), PInches(8.3), PInches(3.2))
                box.text_frame.text = "See PDF/DOCX for full content + exercises."

        # MCQs slide
        mc = prs.slides.add_slide(prs.slide_layouts[1])
        mc.shapes.title.text = f"{ch} — MCQs"
        tfm = mc.placeholders[1].text_frame
        tfm.clear()
        tfm.text = "MCQs (Answer key included):"
        for i,(q,_opts,key) in enumerate(MCQS[ch],1):
            p = tfm.add_paragraph()
            p.text = f"Q{i}. {q} (Ans: {key})"
            p.level = 1

    end = prs.slides.add_slide(prs.slide_layouts[1])
    end.shapes.title.text = "StudentDB Projects"
    tf = end.placeholders[1].text_frame
    tf.clear()
    tf.text = "Included in PDF Appendices:"
    p=tf.add_paragraph(); p.text="• studentdb_pickle_undo_redo.py (Pickle + Undo/Redo)"; p.level=1
    p=tf.add_paragraph(); p.text="• studentdb_json_undo_redo.py (JSON + Undo/Redo)"; p.level=1

    prs.save(PPTX_PATH)

def main():
    # generate diagram images
    diag_paths = {}
    for tid in DIAGRAM_TOPIC_IDS:
        diag_paths[tid] = diagram_for_topic(tid)

    # build exports
    build_pdf(diag_paths)
    build_docx(diag_paths)
    build_ppt(diag_paths)

    # write project .py files too
    with open(os.path.join(OUT_DIR, "studentdb_pickle_undo_redo.py"), "w", encoding="utf-8") as f:
        f.write(STUDENTDB_PICKLE_CODE)
    with open(os.path.join(OUT_DIR, "studentdb_json_undo_redo.py"), "w", encoding="utf-8") as f:
        f.write(STUDENTDB_JSON_CODE)

    print("DONE ✅")
    print("Created:")
    print(" -", PDF_PATH)
    print(" -", DOCX_PATH)
    print(" -", PPTX_PATH)
    print(" - studentdb_pickle_undo_redo.py")
    print(" - studentdb_json_undo_redo.py")

if __name__ == "__main__":
    main()
